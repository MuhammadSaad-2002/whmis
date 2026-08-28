#!/usr/bin/env python3
"""Build the WHMIS executive pitch deck as a 16:9 PowerPoint (.pptx).

Reuses the screenshots in docs/presentation/img/ (captured by
scripts/capture-screenshots.mjs) and mirrors the branding of the Word walkthrough
(scripts/build-presentation.py). Produces
docs/presentation/WHMIS-System-Presentation.pptx.

    python3 scripts/build-pptx.py
"""
import os
from datetime import date

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG = os.path.join(ROOT, "docs", "presentation", "img")
OUT = os.path.join(ROOT, "docs", "presentation", "WHMIS-System-Presentation.pptx")

CLIENT = "MASTER PHARMACUTICALS DISTRIBUTOR"
VENDOR = "Virtual Wisdom Technologies (SMC-Private) Limited"
YEAR = date.today().year

# Palette (matches the docx).
SLATE = RGBColor(0x0F, 0x17, 0x2A)
SLATE_2 = RGBColor(0x1E, 0x29, 0x3B)
EMERALD = RGBColor(0x05, 0x96, 0x69)
EMERALD_DK = RGBColor(0x04, 0x78, 0x57)
EMERALD_LT = RGBColor(0x6E, 0xE7, 0xB7)
GREY = RGBColor(0x6B, 0x72, 0x80)
GREY_LT = RGBColor(0x9C, 0xA3, 0xAF)
INK = RGBColor(0x1F, 0x29, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xD1, 0xD5, 0xDB)
BG = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

# 16:9 canvas.
EMU_IN = 914400
SW = 13.333
SH = 7.5

# --------------------------------------------------------------------------- #
# Content
# --------------------------------------------------------------------------- #
CHALLENGE = {
    "eyebrow": "THE OPPORTUNITY",
    "title": "Pharmaceutical distribution is uniquely unforgiving",
    "tagline": "The everyday details that quietly erode a distributor's margin",
    "points": [
        ("Bonus stock hides your true cost",
         "Free “bonus” goods change what stock really costs — ignore them and every margin looks better than it is."),
        ("Expiry is money on a timer",
         "Batches expire; ship the wrong one and stock becomes a write-off instead of a sale."),
        ("Credit sales strain cash flow",
         "Pharmacies buy on credit — without limits and aging, overdue money hides in plain sight."),
        ("Trade schemes applied by hand drift",
         "Inconsistent, unmeasured schemes cost more than you think and treat customers unequally."),
    ],
    "close": "WHMIS is engineered around exactly these realities — so the system protects your margin while your team simply does its job.",
}

THEMES = [
    {
        "eyebrow": "THE PLATFORM",
        "title": "One connected system",
        "tagline": "Purchase → stock → book → sell → collect → report — reconciled automatically",
        "bullets": [
            "Every step of the business in a single, always-current source of truth — no more scattered spreadsheets and registers.",
            "Purpose-built for pharma: batch & expiry tracking, bonus-stock handling, and a true cost that includes free goods.",
            "Localised for Pakistan — PKR pricing with GST / NTN / STRN / CNIC on every party and document.",
            "Runs on ordinary cPanel hosting, in any browser — nothing special to install.",
        ],
        "img": "02-dashboard.png",
        "note": "The dashboard is the first screen the team sees each day — today's sales and purchases, month sales and profit, receivables vs. payables, inventory value, and a 12-month trend. It turns hindsight into foresight.",
    },
    {
        "eyebrow": "TRANSACTIONS",
        "title": "Sell & book — faster billing, controlled credit",
        "tagline": "From field order to posted invoice, without re-keying or risk",
        "bullets": [
            "Keyboard-fast invoice entry built for high daily volumes — bill in seconds, not minutes.",
            "Stock is drawn from the earliest-expiry batch automatically; you never oversell or ship expired goods.",
            "Customers over their credit limit are blocked before a credit invoice can post — cash flow protected.",
            "Field bookings flow through approval into a ready-to-bill invoice, carrying their trade schemes with them.",
        ],
        "img": "04-sales-entry.png",
        "note": "Sales invoices are raised on a keyboard-driven grid; pricing, discount, GST, batch stock and credit are applied automatically and re-verified on the server before posting. Bookings mirror the same fast entry, then submit → approve → convert to a draft sale in one step.",
    },
    {
        "eyebrow": "INVENTORY",
        "title": "Buy & stock — true cost, zero expiry waste",
        "tagline": "Honest margins, and older stock out the door before it expires",
        "bullets": [
            "Free bonus units fold into each batch's effective cost, so profit is always calculated truthfully.",
            "Expiry-first (FIFO) issuing pushes the earliest-expiring batch out first, minimising write-offs.",
            "Near-expiry and low-stock situations are surfaced early — you act before a loss or a shortage.",
            "Every unit's movement in and out is recorded in an append-only history you can always trust.",
        ],
        "img": "10-batches.png",
        "note": "Purchases capture batch, expiry and bonus units; on posting, stock is created and bonus dilutes the batch's effective cost. Stock is tracked per batch, consumed earliest-expiry-first, with full traceability.",
    },
    {
        "eyebrow": "FINANCE",
        "title": "Money & position — know where you stand",
        "tagline": "Receivables vs. payables, and the net position between them — right now",
        "bullets": [
            "A single Financial Position screen: Due from Customers, Owed to Suppliers, and Net Position, headlined.",
            "Receivables and payables with aging buckets, so overdue accounts and upcoming bills stand out instantly.",
            "Receipts and payments allocate against invoices; balances and aging update the moment money moves.",
            "Pick any date range and print the whole position statement to PDF for owners, banks or auditors.",
        ],
        "img": "24-financial-position.png",
        "note": "Profit on paper means little if cash is stuck in receivables while supplier bills come due. The Financial Position screen answers 'can we meet our commitments this week?' — receivables, payables, net position, and a full payment log for any period.",
    },
    {
        "eyebrow": "SALES ENABLEMENT",
        "title": "Trade schemes & clean returns",
        "tagline": "Consistent pricing you can measure — and reversals that never drift",
        "bullets": [
            "Define each scheme once — bonus, slab bonus, percentage or fixed discount, special price — and target it precisely.",
            "Rules only fill the invoice line, so your team keeps the final say; the Incentives Given report shows their true cost.",
            "Sales returns restore stock to the exact batches sold and refund proportionally; purchase returns reverse supplier balances.",
            "Returns post immediately with the right credit and debit notes — inventory and ledgers stay perfectly in step.",
        ],
        "img": "15-incentives.png",
        "note": "Schemes are a major cost most distributors can't quantify. Defining them centrally makes pricing consistent and, for the first time, measurable. Returns are handled exactly — back to the right batch, at the right value — with no manual adjustment.",
    },
    {
        "eyebrow": "INSIGHTS",
        "title": "Seventeen reports, one click to export",
        "tagline": "The views a distributor needs to run — and defend — the business",
        "bullets": [
            "Sales, Purchases, Inventory and Financial reports — registers, profitability, aging, stock, expiry and more.",
            "Each renders with totals and, where useful, a chart — no spreadsheets to build or maintain.",
            "One click exports any report to Excel or PDF for owners, auditors and partners.",
            "Everyone works from the same trusted figures — you spend time deciding, not assembling.",
        ],
        "img": "19-report-profit.png",
        "note": "Seventeen ready-made reports across Sales, Purchases, Inventory and Financial. Nothing to configure — the insight is already there, on screen and exportable.",
    },
    {
        "eyebrow": "PLATFORM",
        "title": "Modern, safe, and localised",
        "tagline": "Many screens in one window, the right access for every role",
        "bullets": [
            "Open sales, bookings, inventory and reports as tabs in one window — switch instantly, each keeps its own state.",
            "Role-based access gates every sensitive action; bookers create orders but never touch billing or master pricing.",
            "Proactive alerts flag low stock, expiry, overdue balances and pending approvals before they cost you.",
            "Browser-based and cPanel-friendly, with your data in your own database, under your control.",
        ],
        "img": "21-workspace.png",
        "note": "A tabbed workspace keeps the team in flow with no lost work. Role-based permissions protect data and margin; in-app alerts surface issues early. Deployment is low-risk — ordinary hosting, any browser, your data.",
    },
]

WHY_FIT = {
    "eyebrow": "THE FIT",
    "title": f"Why WHMIS fits {CLIENT.title()}",
    "tagline": "Built for how a pharmaceutical distributor in Pakistan actually works",
    "points": [
        "One connected workflow — purchase to collection — with every figure reconciled automatically.",
        "Purpose-built for pharma: batches, expiry, bonus stock, and a true cost that keeps margins honest.",
        "Cash under control — credit limits, aging and a live net position protect working capital.",
        "Trade schemes applied consistently and, at last, measured.",
        "Localised for Pakistan: PKR and GST / NTN / STRN / CNIC throughout.",
        "Low-risk to adopt: ordinary cPanel hosting, browser-based, your data in your own database.",
    ],
}


# --------------------------------------------------------------------------- #
# Low-level helpers
# --------------------------------------------------------------------------- #
def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _no_autofit(tf):
    # Prevent PowerPoint from shrinking our text; we size it deliberately.
    tf.word_wrap = True


def rect(slide, x, y, w, h, color, line=None, line_w=None, shadow=False, radius=None):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    if color is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = color
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w or 0.75)
    s.shadow.inherit = False
    if shadow:
        _drop_shadow(s)
    return s


def _drop_shadow(shape):
    spPr = shape._element.spPr
    effLst = spPr.makeelement(qn("a:effectLst"), {})
    shdw = effLst.makeelement(qn("a:outerShdw"), {
        "blurRad": "90000", "dist": "50000", "dir": "5400000", "rotWithShape": "0",
    })
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": "0F172A"})
    alpha = clr.makeelement(qn("a:alpha"), {"val": "22000"})
    clr.append(alpha)
    shdw.append(clr)
    effLst.append(shdw)
    spPr.append(effLst)


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def run(p, text, size, color, bold=False, italic=False, font=FONT, spacing=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    if spacing is not None:
        _letter_spacing(r, spacing)
    return r


def _letter_spacing(r, pts):
    rPr = r._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def para(tf, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    return p


def add_footer(slide, number):
    line = rect(slide, 0.55, 7.02, 12.23, 0.006, CARD_BORDER)
    tb, tf = textbox(slide, 0.55, 7.06, 10.5, 0.35)
    p = tf.paragraphs[0]
    run(p, "Confidential", 8, EMERALD_DK, bold=True)
    run(p, "   ·   WHMIS   ·   " + VENDOR, 8, GREY)
    tb2, tf2 = textbox(slide, 11.0, 7.06, 1.78, 0.35)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run(p2, str(number), 8, GREY)


def eyebrow_title(slide, eyebrow, title, tagline, x=0.55, y=0.5, w=12.2):
    tb, tf = textbox(slide, x, y, w, 0.35)
    run(tf.paragraphs[0], eyebrow, 11, EMERALD, bold=True, spacing=1.5)
    tb2, tf2 = textbox(slide, x, y + 0.34, w, 0.9)
    run(tf2.paragraphs[0], title, 30, SLATE, bold=True)
    # emerald rule
    rect(slide, x + 0.02, y + 1.18, 0.7, 0.045, EMERALD)
    if tagline:
        tb3, tf3 = textbox(slide, x, y + 1.30, w, 0.5)
        run(tf3.paragraphs[0], tagline, 14, GREY)


def framed_image(slide, name, x, y, max_w, max_h):
    path = os.path.join(IMG, name)
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = iw / ih
    w = max_w
    h = w / ratio
    if h > max_h:
        h = max_h
        w = h * ratio
    # center within the box
    cx = x + (max_w - w) / 2
    cy = y + (max_h - h) / 2
    pad = 0.06
    card = rect(slide, cx - pad, cy - pad, w + 2 * pad, h + 2 * pad, WHITE,
                line=CARD_BORDER, line_w=1.0, shadow=True, radius=0.03)
    slide.shapes.add_picture(path, Inches(cx), Inches(cy), Inches(w), Inches(h))
    return cx, cy, w, h


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    r = rect(slide, -0.06, -0.06, SW + 0.12, SH + 0.12, color)
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #
def slide_title(prs):
    s = blank(prs)
    bg(s, SLATE)
    # emerald side accent
    rect(s, 0, 0, 0.22, SH, EMERALD)
    tb, tf = textbox(s, 1.1, 1.75, 11.0, 2.6)
    run(tf.paragraphs[0], "WHMIS", 66, WHITE, bold=True)
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    run(p, "Pharmaceutical Distribution ERP", 24, EMERALD_LT)
    rect(s, 1.14, 3.95, 2.2, 0.05, EMERALD)
    tb2, tf2 = textbox(s, 1.14, 4.15, 11.0, 0.5)
    run(tf2.paragraphs[0], "System Presentation", 16, GREY_LT)

    tb3, tf3 = textbox(s, 1.14, 5.35, 11.0, 1.5)
    run(tf3.paragraphs[0], "PREPARED FOR", 10, GREY_LT, bold=True, spacing=1.5)
    p = tf3.add_paragraph()
    p.space_before = Pt(2)
    run(p, CLIENT, 20, WHITE, bold=True)
    p = tf3.add_paragraph()
    p.space_before = Pt(10)
    run(p, "A product of ", 11, GREY_LT)
    run(p, VENDOR, 13, EMERALD_LT, bold=True)

    tb4, tf4 = textbox(s, 1.14, 6.85, 11.0, 0.4)
    run(tf4.paragraphs[0], date.today().strftime("%B %Y") + "    ·    Confidential", 10, GREY_LT)
    return s


def slide_challenge(prs, number):
    s = blank(prs)
    bg(s, BG)
    eyebrow_title(s, CHALLENGE["eyebrow"], CHALLENGE["title"], CHALLENGE["tagline"])
    top = 2.35
    card_w = 5.9
    gap_x = 0.4
    card_h = 1.55
    gap_y = 0.28
    positions = [(0.55, top), (0.55 + card_w + gap_x, top),
                 (0.55, top + card_h + gap_y), (0.55 + card_w + gap_x, top + card_h + gap_y)]
    for (lead, desc), (x, y) in zip(CHALLENGE["points"], positions):
        rect(s, x, y, card_w, card_h, RGBColor(0xF9, 0xFA, 0xFB), line=CARD_BORDER, line_w=0.75, radius=0.04)
        rect(s, x, y, 0.07, card_h, EMERALD)
        tb, tf = textbox(s, x + 0.32, y + 0.22, card_w - 0.55, card_h - 0.4, anchor=MSO_ANCHOR.MIDDLE)
        run(tf.paragraphs[0], lead, 14.5, SLATE, bold=True)
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        run(p, desc, 11.5, GREY)
    # closing band
    y = top + 2 * card_h + gap_y + 0.24
    band = rect(s, 0.55, y, 12.23, 0.62, RGBColor(0xEC, 0xFD, 0xF5), line=EMERALD_LT, line_w=0.75, radius=0.06)
    tb, tf = textbox(s, 0.9, y, 11.6, 0.62, anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], CHALLENGE["close"], 12.5, EMERALD_DK, bold=True)
    add_footer(s, number)
    set_notes(s, "Frame the client's world before the product: these four realities are where distributor profit quietly leaks. WHMIS is built around them.")
    return s


def slide_theme(prs, theme, number, image_left=False):
    s = blank(prs)
    bg(s, BG)
    eyebrow_title(s, theme["eyebrow"], theme["title"], theme["tagline"])

    text_w = 5.25
    img_box_w = 6.55
    img_box_h = 4.35
    img_y = 2.5
    if image_left:
        img_x = 0.55
        text_x = 0.55 + img_box_w + 0.45
    else:
        img_x = SW - 0.55 - img_box_w
        text_x = 0.55

    framed_image(s, theme["img"], img_x, img_y, img_box_w, img_box_h)

    tb, tf = textbox(s, text_x, img_y - 0.05, text_w, 4.6, anchor=MSO_ANCHOR.TOP)
    first = True
    for b in theme["bullets"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(11)
        p.line_spacing = 1.05
        run(p, "▸  ", 13, EMERALD, bold=True)
        run(p, b, 13, INK)
    add_footer(s, number)
    set_notes(s, theme["note"])
    return s


def slide_why_fit(prs, number):
    s = blank(prs)
    bg(s, BG)
    eyebrow_title(s, WHY_FIT["eyebrow"], WHY_FIT["title"], WHY_FIT["tagline"])
    top = 2.5
    col_w = 5.95
    row_h = 1.15
    gap_x = 0.33
    xs = [0.55, 0.55 + col_w + gap_x]
    for i, point in enumerate(WHY_FIT["points"]):
        col = i % 2
        row = i // 2
        x = xs[col]
        y = top + row * (row_h + 0.12)
        rect(s, x, y, 0.42, 0.42, EMERALD, radius=0.5)
        tbn, tfn = textbox(s, x, y, 0.42, 0.42, anchor=MSO_ANCHOR.MIDDLE)
        tfn.paragraphs[0].alignment = PP_ALIGN.CENTER
        run(tfn.paragraphs[0], str(i + 1), 13, WHITE, bold=True)
        tb, tf = textbox(s, x + 0.6, y - 0.05, col_w - 0.7, row_h, anchor=MSO_ANCHOR.MIDDLE)
        run(tf.paragraphs[0], point, 13, INK)
    add_footer(s, number)
    set_notes(s, "Close the loop: restate that every pharma-specific control they need is already here on day one — this is their requirement, met.")
    return s


def slide_closing(prs):
    s = blank(prs)
    bg(s, SLATE)
    rect(s, 0, 0, 0.22, SH, EMERALD)
    tb, tf = textbox(s, 1.1, 2.3, 11.0, 2.5)
    run(tf.paragraphs[0], "See it on your own data.", 40, WHITE, bold=True)
    p = tf.add_paragraph()
    p.space_before = Pt(14)
    run(p, "We would be glad to walk the team at ", 15, GREY_LT)
    run(p, CLIENT.title(), 15, EMERALD_LT, bold=True)
    p2 = tf.add_paragraph()
    run(p2, "through a live, hands-on demonstration of WHMIS.", 15, GREY_LT)
    rect(s, 1.14, 4.85, 2.2, 0.05, EMERALD)
    tb2, tf2 = textbox(s, 1.14, 5.1, 11.0, 0.9)
    run(tf2.paragraphs[0], "A product of", 11, GREY_LT)
    p = tf2.add_paragraph()
    run(p, VENDOR, 16, WHITE, bold=True)
    return s


def slide_confidential(prs, number):
    s = blank(prs)
    bg(s, BG)
    eyebrow_title(s, "NOTICE", "Confidentiality", "", w=12.2)
    y = 2.3
    band = rect(s, 0.55, y, 12.23, 3.4, RGBColor(0xF9, 0xFA, 0xFB), line=CARD_BORDER, line_w=0.75, radius=0.03)
    rect(s, 0.55, y, 0.09, 3.4, EMERALD)
    tb, tf = textbox(s, 1.05, y + 0.35, 11.3, 2.8)
    run(tf.paragraphs[0],
        f"This document and the WHMIS platform it describes are the confidential and proprietary "
        f"property of {VENDOR}. It has been prepared exclusively for {CLIENT} for the sole purpose "
        f"of evaluating the system.", 13.5, INK)
    p = tf.add_paragraph()
    p.space_before = Pt(12)
    run(p,
        "Its contents may not be copied, reproduced, distributed or disclosed to any third party, "
        "in whole or in part, without the prior written consent of " + VENDOR + ".", 13.5, INK)
    p = tf.add_paragraph()
    p.space_before = Pt(16)
    run(p, f"© {YEAR} {VENDOR}. All rights reserved.", 11, GREY, italic=True)
    add_footer(s, number)
    return s


def build():
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN))
    prs.slide_height = Emu(int(SH * EMU_IN))

    slide_title(prs)
    n = 2
    slide_challenge(prs, n); n += 1
    for i, theme in enumerate(THEMES):
        slide_theme(prs, theme, n, image_left=(i % 2 == 1)); n += 1
    slide_why_fit(prs, n); n += 1
    slide_closing(prs)
    slide_confidential(prs, n)

    prs.save(OUT)
    print(f"Wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
